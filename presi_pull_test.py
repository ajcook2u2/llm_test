from pptx import Presentation
from pptx.util import Inches, Pt
import pprint
import pandas as pd
from analysis_generator import data_delta, response

#lowest should be first quarter of 2024
year_selection = 2024
quarter_selection = 1


prs = Presentation('E:/LibreOffice/test presi.pptx')

slide = prs.slides[0]
slide1 = prs.slides[1]
slide2 = prs.slides[2]


current_quarter = pd.read_csv(f'{year_selection}-quarterly.csv')
quarter_columns = current_quarter.columns
current_quarter = current_quarter.iloc[quarter_selection - 1]

current_month = pd.read_csv(f'{year_selection}-monthly.csv')
month_columns = current_month.columns
current_month = current_month.iloc[(1 * (quarter_selection - 1)):(3 * (quarter_selection - 1))]

yoy_quarter = pd.read_csv(f'{year_selection - 1}-quarterly.csv')
yoy_quarter = yoy_quarter.iloc[quarter_selection - 1]

yoy_month = pd.read_csv(f'{year_selection - 1}-monthly.csv')
yoy_month = yoy_month.iloc[(1 * (quarter_selection - 1)):(3 * (quarter_selection - 1))]

if quarter_selection == 1:
    qoq_quarter = pd.read_csv(f'{year_selection - 1}-quarterly.csv')
    qoq_quarter = qoq_quarter.iloc[quarter_selection + 2]
else:
    qoq_quarter = pd.read_csv(f'{year_selection}-quarterly.csv')
    qoq_quarter = qoq_quarter.iloc[quarter_selection - 2]


qoq_month = pd.read_csv('2023-monthly.csv')
qoq_month = qoq_month.iloc[-3:]
qoq_month.reset_index(inplace=True)

deltas = data_delta(current_quarter, qoq_quarter, yoy_quarter, current_month, qoq_month, yoy_month, quarter_columns)

#the asset tree of a slide, from top-down order
for i in slide1.shapes:
    print(i)



def number_round_stylized(number):
    number = float(number)
    if 1 <= number < 1000:
        return number
    elif 1000 <= number < 1000000:
        number = number / 10
        number = round(number)
        number = str(number / 100)
        number = f'{number}k'
        return number
    elif number > 1000000:
        number = number / 10000
        number = round(number)
        number = str(number / 100)
        number = f'{number}M'
        return number
    else:
        number1 = str(number)
        number1 = number1.split('.')[1]
        scaler = 2
        for digit in number1:
            if digit == '0':
                scaler += 1
            else:
                break
        if scaler > 0:
            number = number * (10**scaler)
            number = round(number)
            number = number / (10**scaler)
            return number
        else:
            number = float(number)
            number = number * 1000
            number = round(number)
            number = number / 1000
            return number


#how to access the title of a slide
title = slide1.shapes.title
print('title', title.text)
print(current_quarter['cpc'], 'k')
print(number_round_stylized(current_quarter['cpc']), 'p')

#How to access and change cells on a table
table = slide1.shapes[1]
print(table.table.cell(0, 0).text)
table.table.cell(0, 0).text = 'Year'
table.table.cell(1, 0).text = '2024'
table.table.cell(2, 0).text = '2023'
table.table.cell(3, 0).text = 'YoY'
#
table.table.cell(0, 1).text = 'Spend'
table.table.cell(1, 1).text = f"${number_round_stylized(current_quarter['spend'])}"
table.table.cell(2, 1).text = f"${number_round_stylized(yoy_quarter['spend'])}"
table.table.cell(3, 1).text = f"{deltas['yoy_quarter_delta'][0]['spend']}"
#
table.table.cell(0, 2).text = 'Impr.'
table.table.cell(1, 2).text = f"{number_round_stylized(current_quarter['impr'])}"
table.table.cell(2, 2).text = f"{number_round_stylized(yoy_quarter['impr'])}"
table.table.cell(3, 2).text = f"{deltas['yoy_quarter_delta'][1]['impr']}"

table.table.cell(0, 3).text = 'Clicks'
table.table.cell(1, 3).text = f"{number_round_stylized(current_quarter['clicks'])}"
table.table.cell(2, 3).text = f"{number_round_stylized(yoy_quarter['clicks'])}"
table.table.cell(3, 3).text = f"{deltas['yoy_quarter_delta'][2]['clicks']}"

table.table.cell(0, 4).text = 'CTR'
table.table.cell(1, 4).text = f"{number_round_stylized(current_quarter['ctr'])}%"
table.table.cell(2, 4).text = f"{number_round_stylized(yoy_quarter['ctr'])}%"
table.table.cell(3, 4).text = f"{deltas['yoy_quarter_delta'][3]['ctr']}"

table.table.cell(0, 5).text = 'CPC'
table.table.cell(1, 5).text = f"${number_round_stylized(current_quarter['cpc'])}"
table.table.cell(2, 5).text = f"${number_round_stylized(yoy_quarter['cpc'])}"
table.table.cell(3, 5).text = f"{deltas['yoy_quarter_delta'][6]['cpc']}"

table.table.cell(0, 6).text = 'Conversions'
table.table.cell(1, 6).text = f"{number_round_stylized(current_quarter['conversions'])}"
table.table.cell(2, 6).text = f"{number_round_stylized(yoy_quarter['conversions'])}"
table.table.cell(3, 6).text = f"{deltas['yoy_quarter_delta'][4]['conversions']}"

table.table.cell(0, 7).text = 'Conv. Rate'
table.table.cell(1, 7).text = f"{number_round_stylized(current_quarter['conv rate'])}"
table.table.cell(2, 7).text = f"{number_round_stylized(yoy_quarter['conv rate'])}"
table.table.cell(3, 7).text = f"{deltas['yoy_quarter_delta'][7]['conv rate']}"



#how to access the text box,
text = slide1.shapes[2]
print('text', text.text)

prs.save('pull_test.pptx')
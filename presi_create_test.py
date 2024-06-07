from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
blank_slide_layout = prs.slide_layouts[6]

# Create two separate slides:
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]  # Access the first placeholder (title) and second one for subtitle
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

slide2 = prs.slides.add_slide(blank_slide_layout)

width = height = Inches(1)
left = Inches(1.5)
top = Inches(-.25)
title_box = slide2.shapes.add_textbox(left, top, width, height)
title_box = title_box.text_frame

title = title_box.add_paragraph()
title.text = 'Lorem Ipsum'
title.font.size = Pt(40)

width = height = Inches(5)
left = Inches(1.5)
top = Inches(5)
text_box = slide2.shapes.add_textbox(left, top, width, height)
text_box = text_box.text_frame

text = text_box.add_paragraph()
text.text = 'Lorem Ipsum12'
text.font.size = Pt(12)

prs.save('test.pptx')